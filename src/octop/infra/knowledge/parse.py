"""Text extraction for the document types accepted by knowledge bases."""

from __future__ import annotations

import csv
import io
import json
from collections.abc import Iterable
from html.parser import HTMLParser
from pathlib import Path
from typing import TYPE_CHECKING, Any

from octop.infra.knowledge.ocr import OCR_IMAGE_SUFFIXES

if TYPE_CHECKING:
    from octop.infra.knowledge.ocr import OcrExtractor

_PLAIN_TEXT_SUFFIXES = {
    ".md",
    ".markdown",
    ".txt",
    ".rst",
    ".yaml",
    ".yml",
    ".jsonl",
}
_DOCX_FALLBACK_TAG = "{http://schemas.openxmlformats.org/markup-compatibility/2006}Fallback"
_DOCX_HTML_TYPES = {"application/xhtml+xml", "text/html"}
_DOCX_MAIN_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def parse_document(path: Path, *, ocr: OcrExtractor | None = None) -> str:
    """Extract searchable text from a supported local document."""
    suffix = path.suffix.lower()
    if suffix in OCR_IMAGE_SUFFIXES:
        if ocr is None:
            raise RuntimeError("knowledge OCR is not enabled")
        return ocr(path)
    if suffix in _PLAIN_TEXT_SUFFIXES:
        return _read_text(path)
    if suffix == ".json":
        return _parse_json(path)
    if suffix in {".html", ".htm"}:
        return _parse_html(path)
    if suffix == ".csv":
        return _parse_delimited(path, delimiter=",")
    if suffix == ".tsv":
        return _parse_delimited(path, delimiter="\t")
    if suffix == ".pdf":
        from pypdf import PdfReader

        text = "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
        if text.strip() or ocr is None:
            return text
        return ocr(path)
    if suffix == ".docx":
        return _parse_docx(path)
    if suffix == ".pptx":
        from pptx import Presentation

        return "\n".join(
            shape.text
            for slide in Presentation(str(path)).slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
    if suffix in {".xlsx", ".xlsm"}:
        return _parse_xlsx(path)
    if suffix == ".xls":
        return _parse_xls(path)
    raise ValueError(f"unsupported knowledge document extension: {suffix or '(none)'}")


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8-sig", errors="replace")


def _parse_json(path: Path) -> str:
    raw = _read_text(path)
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        return raw
    return json.dumps(parsed, ensure_ascii=False, indent=2)


def _parse_docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    return _docx_body_text(document.element.body, document.part)


def _docx_body_text(body: Any, part: Any) -> str:
    """Collect every ``w:p`` and ``w:altChunk`` in a body, in document order.

    ``Document.paragraphs`` only lists ``w:p`` children of ``w:body``, so text inside
    tables, content controls (``w:sdt``), revision wrappers (``w:ins``), and text boxes
    is silently dropped. ``mc:Fallback`` duplicates its ``mc:Choice`` sibling, so its
    paragraphs are skipped.
    """
    from docx.oxml.ns import qn

    alt_chunk_tag = qn("w:altChunk")
    lines: list[str] = []
    for element in body.iter(qn("w:p"), alt_chunk_tag):
        if next(element.iterancestors(_DOCX_FALLBACK_TAG), None) is not None:
            continue
        if element.tag == alt_chunk_tag:
            lines.append(_docx_alt_chunk_text(element, part))
        else:
            lines.append(element.text)
    return "\n".join(lines)


def _docx_alt_chunk_text(element: Any, part: Any) -> str:
    """Expand a ``w:altChunk`` reference the way Word does when opening the file.

    Converters keep the bulk of the text in an embedded HTML or Word part and leave
    only a stub in ``document.xml``; unexpanded, that body is lost.
    """
    from docx.oxml.ns import qn

    chunk = part.related_parts.get(element.get(qn("r:id")) or "")
    if chunk is None:
        return ""
    content_type = str(chunk.content_type)
    blob: bytes = chunk.blob
    if content_type in _DOCX_HTML_TYPES:
        return _html_text(blob.decode("utf-8-sig", errors="replace"))
    if content_type == "text/plain":
        return blob.decode("utf-8-sig", errors="replace")
    if content_type == _DOCX_MAIN_TYPE:
        from docx import Document

        nested = Document(io.BytesIO(blob))
        return _docx_body_text(nested.element.body, nested.part)
    return ""


def _parse_html(path: Path) -> str:
    return _html_text(_read_text(path))


def _html_text(raw: str) -> str:
    parser = _HTMLTextParser()
    parser.feed(raw)
    parser.close()
    lines = [" ".join(line.split()) for line in parser.text().splitlines()]
    return "\n".join(line for line in lines if line)


class _HTMLTextParser(HTMLParser):
    _SKIP = frozenset({"script", "style", "noscript", "template"})
    _BLOCK = frozenset(
        {
            "p",
            "div",
            "br",
            "li",
            "tr",
            "h1",
            "h2",
            "h3",
            "h4",
            "h5",
            "h6",
            "article",
            "section",
            "header",
            "footer",
            "blockquote",
            "pre",
            "table",
            "ul",
            "ol",
        }
    )

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._parts: list[str] = []
        self._skip = 0

    def handle_starttag(self, tag: str, _attrs: list[tuple[str, str | None]]) -> None:
        if tag in self._SKIP:
            self._skip += 1
            return
        if not self._skip and tag in self._BLOCK:
            self._parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP:
            if self._skip:
                self._skip -= 1
            return
        if not self._skip and tag in self._BLOCK:
            self._parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._skip and data:
            self._parts.append(data)

    def text(self) -> str:
        return "".join(self._parts)


def _stringify_cell(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).replace("\t", " ").replace("\n", " ").strip()


def _trim_row(cells: list[str]) -> list[str]:
    while cells and cells[-1] == "":
        cells.pop()
    return cells


def _sheet_text(title: str, rows: Iterable[Iterable[object]]) -> str:
    lines = [f"# {title}"]
    has_cells = False
    for raw in rows:
        cells = _trim_row([_stringify_cell(cell) for cell in raw])
        if not cells:
            continue
        has_cells = True
        lines.append("\t".join(cells))
    if not has_cells:
        return ""
    return "\n".join(lines)


def _parse_delimited(path: Path, *, delimiter: str) -> str:
    reader = csv.reader(_read_text(path).splitlines(), delimiter=delimiter)
    return _sheet_text(path.stem, reader)


def _parse_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for sheet in workbook.worksheets:
            text = _sheet_text(str(sheet.title), sheet.iter_rows(values_only=True))
            if text:
                parts.append(text)
        return "\n\n".join(parts)
    finally:
        workbook.close()


def _parse_xls(path: Path) -> str:
    import xlrd

    book = xlrd.open_workbook(str(path), formatting_info=False)
    parts: list[str] = []
    for sheet in book.sheets():
        rows: list[list[object]] = []
        for row_idx in range(sheet.nrows):
            rows.append(
                [
                    _xlrd_cell_value(book, sheet.cell(row_idx, col_idx))
                    for col_idx in range(sheet.ncols)
                ]
            )
        text = _sheet_text(sheet.name, rows)
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def _xlrd_cell_value(book: object, cell: object) -> object:
    import xlrd

    ctype = getattr(cell, "ctype", xlrd.XL_CELL_EMPTY)
    value = getattr(cell, "value", None)
    if ctype in {xlrd.XL_CELL_EMPTY, xlrd.XL_CELL_BLANK, xlrd.XL_CELL_ERROR}:
        return None
    if ctype == xlrd.XL_CELL_DATE:
        datemode = int(getattr(book, "datemode", 0))
        try:
            return xlrd.xldate_as_datetime(value, datemode).isoformat(sep=" ", timespec="seconds")
        except (OSError, OverflowError, TypeError, ValueError):
            return value
    return value
