"""Unit tests for supported knowledge-document parsers."""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest

from octop.infra.knowledge.parse import parse_document

_DOCX_NAMESPACES = (
    'xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
    'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
    'xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape" '
    'xmlns:v="urn:schemas-microsoft-com:vml"'
)
_DOCX_MAIN_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_ALT_CHUNK_RELATIONSHIP = (
    '<Relationship Id="rId99" Target="../{name}" '
    'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/aFChunk"/>'
)


def _docx_paragraph(text: str) -> str:
    return f"<w:p><w:r><w:t>{text}</w:t></w:r></w:p>"


def _docx_with_body(
    path: Path, body_xml: str, *, alt_chunk: tuple[str, str, bytes] | None = None
) -> Path:
    """Write a .docx whose ``w:body`` holds *body_xml* verbatim.

    *alt_chunk* is ``(part_name, content_type, blob)`` for a part related as ``rId99``,
    matching what converters emit for ``<w:altChunk r:id="rId99"/>``.
    """
    from docx import Document

    Document().save(path)
    document = f"<w:document {_DOCX_NAMESPACES}><w:body>{body_xml}</w:body></w:document>".encode()
    with zipfile.ZipFile(path) as source:
        entries = [(info.filename, source.read(info.filename)) for info in source.infolist()]
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as target:
        for name, data in entries:
            if name == "word/document.xml":
                data = document
            elif alt_chunk is not None and name == "word/_rels/document.xml.rels":
                relationship = _ALT_CHUNK_RELATIONSHIP.format(name=alt_chunk[0])
                data = data.replace(
                    b"</Relationships>", relationship.encode() + b"</Relationships>"
                )
            elif alt_chunk is not None and name == "[Content_Types].xml":
                override = f'<Override PartName="/{alt_chunk[0]}" ContentType="{alt_chunk[1]}"/>'
                data = data.replace(b"</Types>", override.encode() + b"</Types>")
            target.writestr(name, data)
        if alt_chunk is not None:
            target.writestr(alt_chunk[0], alt_chunk[2])
    path.write_bytes(buffer.getvalue())
    return path


def test_parse_plain_text_and_markdown(tmp_path: Path) -> None:
    text = tmp_path / "notes.txt"
    markdown = tmp_path / "readme.md"
    text.write_text("plain notes", encoding="utf-8")
    markdown.write_text("# Heading\n\nbody", encoding="utf-8")

    assert parse_document(text) == "plain notes"
    assert parse_document(markdown) == "# Heading\n\nbody"


def test_parse_pdf_docx_and_pptx(tmp_path: Path) -> None:
    from docx import Document
    from pptx import Presentation
    from pypdf import PdfWriter

    pdf = tmp_path / "empty.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    with pdf.open("wb") as output:
        writer.write(output)

    docx = tmp_path / "notes.docx"
    document = Document()
    document.add_paragraph("Word notes")
    document.save(docx)

    pptx = tmp_path / "slides.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[0]).shapes.title.text = "Slide title"
    presentation.save(pptx)

    assert parse_document(pdf) == ""
    assert parse_document(docx) == "Word notes"
    assert parse_document(pptx) == "Slide title"


def test_parse_docx_reads_tables_content_controls_and_revisions(tmp_path: Path) -> None:
    body = (
        _docx_paragraph("发布说明")
        + "<w:p/>"
        + f"<w:tbl><w:tr><w:tc>{_docx_paragraph('第一章 总则')}</w:tc></w:tr></w:tbl>"
        + f"<w:sdt><w:sdtContent>{_docx_paragraph('第一条 为了规范')}</w:sdtContent></w:sdt>"
        + '<w:ins w:id="1" w:author="u" w:date="2026-08-27T00:00:00Z">'
        + f"{_docx_paragraph('第二条 本办法适用于')}</w:ins>"
    )
    path = _docx_with_body(tmp_path / "nested.docx", body)

    assert parse_document(path) == "发布说明\n\n第一章 总则\n第一条 为了规范\n第二条 本办法适用于"


def test_parse_docx_reads_text_box_once(tmp_path: Path) -> None:
    text_box = f"<w:txbxContent>{_docx_paragraph('文本框内容')}</w:txbxContent>"
    body = _docx_paragraph("正文段落") + (
        "<w:p><w:r><mc:AlternateContent>"
        f'<mc:Choice Requires="wps"><wps:txbx>{text_box}</wps:txbx></mc:Choice>'
        f"<mc:Fallback><v:textbox>{text_box}</v:textbox></mc:Fallback>"
        "</mc:AlternateContent></w:r></w:p>"
    )
    path = _docx_with_body(tmp_path / "textbox.docx", body)

    assert parse_document(path) == "正文段落\n\n文本框内容"


def test_parse_docx_expands_html_alt_chunk(tmp_path: Path) -> None:
    chunk = "<html><body><p>第一章 总则</p><p>第一条 为了规范</p></body></html>".encode()
    path = _docx_with_body(
        tmp_path / "html_chunk.docx",
        _docx_paragraph("发布说明") + '<w:altChunk r:id="rId99"/>',
        alt_chunk=("chunk.xhtml", "application/xhtml+xml", chunk),
    )

    assert parse_document(path) == "发布说明\n第一章 总则\n第一条 为了规范"


def test_parse_docx_expands_word_alt_chunk(tmp_path: Path) -> None:
    nested = _docx_with_body(tmp_path / "nested.docx", _docx_paragraph("嵌套正文"))
    path = _docx_with_body(
        tmp_path / "word_chunk.docx",
        _docx_paragraph("发布说明") + '<w:altChunk r:id="rId99"/>',
        alt_chunk=("chunk.docx", _DOCX_MAIN_TYPE, nested.read_bytes()),
    )

    assert parse_document(path) == "发布说明\n嵌套正文"


def test_parse_csv_xlsx_and_xls(tmp_path: Path) -> None:
    import xlwt
    from openpyxl import Workbook

    csv_path = tmp_path / "sales.csv"
    csv_path.write_text("item,qty\napple,2\n", encoding="utf-8")

    tsv_path = tmp_path / "sales.tsv"
    tsv_path.write_text("item\tqty\napple\t2\n", encoding="utf-8")

    xlsx = tmp_path / "sales.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Q1"
    sheet.append(["item", "qty"])
    sheet.append(["apple", 2])
    extra = workbook.create_sheet("Q2")
    extra.append(["item", "qty"])
    extra.append(["pear", 3])
    workbook.save(xlsx)

    xlsm = tmp_path / "sales.xlsm"
    workbook.save(xlsm)

    xls = tmp_path / "sales.xls"
    book = xlwt.Workbook()
    q1 = book.add_sheet("Q1")
    q1.write(0, 0, "item")
    q1.write(0, 1, "qty")
    q1.write(1, 0, "apple")
    q1.write(1, 1, 2)
    book.save(str(xls))

    expected_xlsx = "# Q1\nitem\tqty\napple\t2\n\n# Q2\nitem\tqty\npear\t3"
    assert parse_document(csv_path) == "# sales\nitem\tqty\napple\t2"
    assert parse_document(tsv_path) == "# sales\nitem\tqty\napple\t2"
    assert parse_document(xlsx) == expected_xlsx
    assert parse_document(xlsm) == expected_xlsx
    assert parse_document(xls) == "# Q1\nitem\tqty\napple\t2"


def test_parse_html_json_and_plain_variants(tmp_path: Path) -> None:
    html = tmp_path / "page.html"
    html.write_text(
        "<html><head><style>p{color:red}</style></head>"
        "<body><h1>Title</h1><p>Hello <b>world</b></p>"
        "<script>alert(1)</script></body></html>",
        encoding="utf-8",
    )
    markdown = tmp_path / "notes.markdown"
    markdown.write_text("# Heading\n\nbody", encoding="utf-8")
    rst = tmp_path / "notes.rst"
    rst.write_text("Heading\n=======\n\nbody", encoding="utf-8")
    yaml_path = tmp_path / "config.yaml"
    yaml_path.write_text("name: octop\n", encoding="utf-8")
    json_path = tmp_path / "data.json"
    json_path.write_text('{"name":"octop","ok":true}', encoding="utf-8")
    jsonl = tmp_path / "rows.jsonl"
    jsonl.write_text('{"a":1}\n{"b":2}\n', encoding="utf-8")

    assert parse_document(html) == "Title\nHello world"
    assert parse_document(markdown) == "# Heading\n\nbody"
    assert parse_document(rst) == "Heading\n=======\n\nbody"
    assert parse_document(yaml_path) == "name: octop\n"
    assert parse_document(json_path) == '{\n  "name": "octop",\n  "ok": true\n}'
    assert parse_document(jsonl) == '{"a":1}\n{"b":2}\n'


def test_parse_rejects_unsupported_extension(tmp_path: Path) -> None:
    path = tmp_path / "unsupported.zip"
    path.write_bytes(b"not a document")

    with pytest.raises(ValueError, match="unsupported"):
        parse_document(path)
